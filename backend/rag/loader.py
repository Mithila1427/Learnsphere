from PyPDF2 import PdfReader
from PyPDF2.errors import PdfReadError
from pptx import Presentation
import os

def load_text(file_path: str) -> str:
    text = ""
    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext == ".pdf":
            reader = PdfReader(file_path)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"

        elif ext == ".pptx":
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

        return text

    except Exception:
        return ""
